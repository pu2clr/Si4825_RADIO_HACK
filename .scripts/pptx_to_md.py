#!/usr/bin/env python3
from pptx import Presentation
from pathlib import Path
import sys
import re


def clean_text(text: str) -> str:
    if not text:
        return ''
    # remove nulls and normalize whitespace
    text = text.replace('\x00', '')
    # replace stray tabs used by pandoc conversion for ranges with en-dash
    text = text.replace('\t', '–')
    # normalize CRLF and multiple newlines
    text = text.replace('\r', '\n')
    text = re.sub('\n{2,}', '\n', text)
    text = text.strip()
    return text


def save_image(shape, out_path: Path, prefix: str, idx: int) -> str:
    image = shape.image
    ext = image.ext
    name = f"{prefix}_img_{idx}.{ext}"
    out_file = out_path / name
    out_file.write_bytes(image.blob)
    return name


def pptx_to_markdown(pptx_path: Path, output_md: Path):
    out_dir = output_md.parent
    images_dir = out_dir / 'images'
    images_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    lines = []
    lines.append(f"# {pptx_path.stem}\n")

    img_counter = 1
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {i}\n")

        # title if present
        title = None
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = clean_text(slide.shapes.title.text)
        except Exception:
            title = None
        if title:
            lines.append(f"### {title}\n")

        # collect text from shapes
        for shape in slide.shapes:
            # pictures
            if getattr(shape, 'shape_type', None) is not None and hasattr(shape, 'image') and shape.shape_type == 13:
                try:
                    saved = save_image(shape, images_dir, f"slide{i}", img_counter)
                    lines.append(f"![Slide {i} image]({output_md.parent.name}/images/{saved})\n")
                    img_counter += 1
                except Exception:
                    pass

            # text-containing shapes
            if hasattr(shape, 'text') and shape.text:
                text = clean_text(shape.text)
                if not text:
                    continue
                # skip duplicate title
                if title and text == title:
                    continue
                # prefix multiple lines as bullets when appropriate
                parts = [p.strip() for p in text.split('\n') if p.strip()]
                if len(parts) == 1:
                    lines.append(parts[0] + "\n")
                else:
                    for p in parts:
                        lines.append(f"- {p}\n")

        lines.append('---\n')

    output_md.parent.mkdir(parents=True, exist_ok=True)
    output_md.write_text("\n".join(lines), encoding='utf-8')
    print(output_md)


if __name__ == '__main__':
    pptx_path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("./Modificando de Si4825 para Si4827 com interface Digital.pptx")
    output = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("./.TransformationArtifacts/si4827_conversion.md")
    pptx_to_markdown(pptx_path, output)
